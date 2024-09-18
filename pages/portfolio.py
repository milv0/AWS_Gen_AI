import PyPDF2, pptx, re, io, boto3, base64
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.docstore.document import Document
from langchain.prompts import PromptTemplate
from langchain.chains.summarize import load_summarize_chain
from pptx.enum.shapes import MSO_SHAPE_TYPE
import streamlit as st
from pages.chatbot import init_conversationchain

# 세션 상태 초기화 Default 값
if "TEMPERATURE" not in st.session_state:
    st.session_state.TEMPERATURE = 0.8
if "TOP_P" not in st.session_state:
    st.session_state.TOP_P = 1.0
if "TOP_K" not in st.session_state:
    st.session_state.TOP_K = 500
if "MAX_TOKENS" not in st.session_state:
    st.session_state.MAX_TOKENS = 4096
if "MEMORY_WINDOW" not in st.session_state:
    st.session_state.MEMORY_WINDOW = 10


def get_portfolio(texts, llm, sections):
    # check korean
    pattern_hangul = re.compile('[\u3131-\u3163\uac00-\ud7a3]+') 
    word_kor = pattern_hangul.search(str(texts))

    if word_kor:
        prompt_template = f"""\n\nHuman: 다음 파일을 분석해서 프로젝트 보고서 형식으로 전환해주세요. ({', '.join(sections)}) 으로 정리해줘. 만약 해당 부분이 없다면 새로 생성해줘:

        {{text}}

        Assistant:"""        
    else:         
        prompt_template = f"""\n\nHuman: Please analyze the following files and convert it into a project report format with the following sections: ({', '.join(sections)}). If any section is missing, please create it.

        {{text}}

        Assistant:"""

    PROMPT = PromptTemplate(template=prompt_template, input_variables=["text"])

    chain = load_summarize_chain(llm, chain_type="stuff", prompt=PROMPT)

    docs = [
        Document(
            page_content=t
        ) for t in texts[:3]
    ]
    summary = chain.invoke(docs)
    print('\n' + summary['output_text'])
    if summary == '':  # error notification
        summary = 'Fail to analysis the document. Try again...'
        return summary
    else:
        return summary


def portfolio():
    # 파일 업로드 또는 텍스트 입력
    upload_or_text = st.radio("텍스트 입력 또는 파일 업로드", ["텍스트 입력", "파일 업로드"])
    text = None  # text 변수를 None으로 초기화

    # 기본 섹션 설정
    all_sections = ["제목", "개요", "배경", "맡은 역할", "결과", "결론", "향후 계획"]

    default_sections = ["제목", "개요", "배경", "맡은 역할", "결과"]
    
    sections = st.multiselect("보고서 섹션을 선택하세요", all_sections, default_sections)

    # 중복 제거
    sections = list(set(sections))
    
    if upload_or_text == "텍스트 입력":
        text_input = st.text_area("텍스트를 입력하세요", height=200)
        if st.button("분석 시작"):
            text = text_input

    elif upload_or_text == "파일 업로드":
        uploaded_file = st.file_uploader("PDF, 텍스트 또는 PowerPoint 파일을 선택하세요", type=["pdf", "txt", "ppt", "pptx"])

        if uploaded_file is not None:
            try:
                if uploaded_file.name.endswith(".pdf"):
                    pdf_reader = PyPDF2.PdfReader(uploaded_file)
                    text = ""
                    for page in range(len(pdf_reader.pages)):
                        text += pdf_reader.pages[page].extract_text()
                elif uploaded_file.name.endswith((".ppt", ".pptx")):
                    presentation = pptx.Presentation(uploaded_file)
                    text = []
                    for slide in presentation.slides:
                        slide_text = []
                        for shape in slide.shapes:
                            if hasattr(shape, "text") and shape.text:
                                slide_text.append(shape.text)
                        text.append("\n".join(slide_text))
                    text = "\n\n".join(text)
                else:
                    text = uploaded_file.read().decode("utf-8")
            except Exception as e:
                st.error(f"파일 읽기 오류: {str(e)}")
                return


    if text:
        text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=100)
        texts = text_splitter.split_text(text)

        summary = get_portfolio(texts, st.session_state.llm, sections)
        st.write(summary['output_text'])

if __name__ == "__main__":
    # 대화 체인 초기화, llm 객체를 session_state에 저장
    conv_chain, llm = init_conversationchain()
    st.session_state.llm = llm

    st.set_page_config(page_title='📝 문서 정리', layout='wide')
    st.title("📝 포트폴리오 정리")

    portfolio()
