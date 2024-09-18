import PyPDF2, pptx, re, io, boto3
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.docstore.document import Document
from langchain.prompts import PromptTemplate
from langchain.chains.summarize import load_summarize_chain
from pptx.enum.shapes import MSO_SHAPE_TYPE
import streamlit as st
from pages.chatbot import init_conversationchain

# 세션 상태 초기화 Default 값
if "TEMPERATURE" not in st.session_state:
    st.session_state.TEMPERATURE = 0.8
if "TOP_P" not in st.session_state:
    st.session_state.TOP_P = 1.0
if "TOP_K" not in st.session_state:
    st.session_state.TOP_K = 500
if "MAX_TOKENS" not in st.session_state:
    st.session_state.MAX_TOKENS = 4096
if "MEMORY_WINDOW" not in st.session_state:
    st.session_state.MEMORY_WINDOW = 10

def get_summary(texts, llm, summary_length):
    # check korean
    pattern_hangul = re.compile('[\u3131-\u3163\uac00-\ud7a3]+') 
    word_kor = pattern_hangul.search(str(texts))

    if word_kor:
        prompt_template = f"""\n\nHuman: 다음 텍스트를 요약해서 {summary_length}자 이내로 설명하세요.

        {{text}}

        Assistant:"""        
    else:         
        prompt_template = f"""\n\nHuman: Write a concise summary of the following text in {summary_length} characters or less:

        {{text}}

        Assistant:"""

    PROMPT = PromptTemplate(template=prompt_template, input_variables=["text"])

    chain = load_summarize_chain(llm, chain_type="stuff", prompt=PROMPT)

    docs = [
        Document(
            page_content=t
        ) for t in texts[:3]
    ]
    summary = chain.invoke(docs)
    print('\n' + summary['output_text'])
    if summary == '':  # error notification
        summary = 'Fail to summarize the document. Try again...'
        return summary
    else:
        return summary

def summarize_document():
    summary_option = st.radio("문서 처리 옵션", ("요약", "정리"))

    # 사용자가 요약 길이를 선택 (요약 옵션일 때만 사용)
    if summary_option == "요약":
        summary_length = st.slider("요약 길이 (자)", min_value=500, max_value=3000, value=2000, step=100)
    else:
        summary_length = None

    # 파일 업로드
    uploaded_file = st.file_uploader("PDF, 텍스트 또는 PowerPoint 파일을 선택하세요", type=["pdf", "txt", "ppt", "pptx"])

    if uploaded_file is not None:
        try:
            if uploaded_file.name.endswith(".pdf"):
                pdf_reader = PyPDF2.PdfReader(uploaded_file)
                text = ""
                for page in range(len(pdf_reader.pages)):
                    text += pdf_reader.pages[page].extract_text()
            elif uploaded_file.name.endswith((".ppt", ".pptx")):
                presentation = pptx.Presentation(uploaded_file)
                text = []
                for slide in presentation.slides:
                    slide_text = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text:
                            slide_text.append(shape.text)
                    text.append("\n".join(slide_text))
                text = "\n\n".join(text)
            else:
                text = uploaded_file.read().decode("utf-8")
        except Exception as e:
            st.error(f"파일 읽기 오류: {str(e)}")
            return

        text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=100)
        texts = text_splitter.split_text(text)

        if summary_option == "요약":
            prompt_template = f"""\n\nHuman: 다음 텍스트를 요약해서 {summary_length}자 이내로 설명하세요.

            {{text}}

            Assistant:"""
        else:
            prompt_template = f"""\n\nHuman: 다음 텍스트 내용을 빠짐 없이 정리해서 설명하세요.

            {{text}}

            Assistant:"""

        PROMPT = PromptTemplate(template=prompt_template, input_variables=["text"])
        chain = load_summarize_chain(st.session_state.llm, chain_type="stuff", prompt=PROMPT)

        docs = [
            Document(
                page_content=t
            ) for t in texts[:3]
        ]
        result = chain.invoke(docs)
        # st.write(result['output_text'])
        st.write(result)

if __name__ == "__main__":
    # 대화 체인 초기화, llm 객체를 session_state에 저장
    conv_chain, llm = init_conversationchain()
    st.session_state.llm = llm    

    st.set_page_config(page_title='📝 문서 처리', layout='wide')
    st.title("📝 문서 처리")

    summarize_document()
